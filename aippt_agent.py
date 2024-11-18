import autogen
import requests
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import time

# OpenRouter API configuration
OPENROUTER_API_KEY = os.environ['OPENROUTER_API_KEY']

class OpenRouterLLM:
    def __init__(self, model="anthropic/claude-3.5-sonnet"):
        self.model = model
        
    def create_completion(self, messages):
        response = requests.post(
            url="https://openrouter.ai/api/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {OPENROUTER_API_KEY}"
            },
            data=json.dumps({
                "model": self.model,
                "messages": messages,
                "max_tokens": 4000,
                "temperature": 0.7,
            })
        )
        return response.json()

class PresentationBuilder:
    def __init__(self, filename="ai_trends_2024.pptx"):
        self.prs = Presentation()
        self.filename = filename
        self.slide_count = 0
        self.target_slides = 15
        
        # Set default slide size to widescreen
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        # Define theme colors
        self.colors = {
            'primary': RGBColor(33, 60, 114),
            'secondary': RGBColor(0, 133, 202),
            'accent': RGBColor(242, 144, 0),
            'text': RGBColor(51, 51, 51),
            'background': RGBColor(255, 255, 255)
        }
        
    def create_title_slide(self, title, subtitle=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        self.slide_count += 1
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['background']
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.color.rgb = self.colors['primary']
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_frame = subtitle_shape.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = self.colors['secondary']
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    def create_content_slide(self, content):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self.slide_count += 1
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['background']
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = content['title']
        title_frame = title_shape.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.color.rgb = self.colors['primary']
        title_para.font.bold = True
        
        # Add body content
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()
        
        for i, bullet in enumerate(content['body']):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = bullet
            p.font.size = Pt(24)
            p.font.color.rgb = self.colors['text']
            p.level = 0
        
        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = content['notes']
        
    def create_chart_slide(self, title, chart_type, chart_data):
        """
        Create a slide with a chart.
        
        Args:
            title (str): The slide title
            chart_type (str): Type of chart ("bar", "line", "pie")
            chart_data (dict): Dictionary containing labels and values arrays
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_count += 1
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.color.rgb = self.colors['primary']
        title_para.font.bold = True
        
        # Create chart placeholder
        chart_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1.5),
            Inches(11),
            Inches(5)
        )
        
        # Style the placeholder
        chart_box.fill.solid()
        chart_box.fill.fore_color.rgb = self.colors['secondary']
        
        # Add data labels
        data_box = slide.shapes.add_textbox(
            Inches(1),
            Inches(1.5),
            Inches(11),
            Inches(5)
        )
        
        # Add data text
        text_frame = data_box.text_frame
        text_frame.clear()
        
        # Add chart type indicator
        p = text_frame.add_paragraph()
        p.text = f"Chart Type: {chart_type.upper()}"
        p.font.size = Pt(14)
        p.font.color.rgb = self.colors['background']
        p.alignment = PP_ALIGN.CENTER
        
        # Add data values
        for label, value in zip(chart_data['labels'], chart_data['values']):
            p = text_frame.add_paragraph()
            p.text = f"{label}: {value}"
            p.font.size = Pt(12)
            p.font.color.rgb = self.colors['background']
            p.alignment = PP_ALIGN.CENTER
        
        # Add caption
        caption_box = slide.shapes.add_textbox(
            Inches(1),
            Inches(6.5),
            Inches(11),
            Inches(0.5)
        )
        caption_frame = caption_box.text_frame
        caption_frame.text = f"Data visualization: {chart_type} chart"
        caption_para = caption_frame.paragraphs[0]
        caption_para.font.size = Pt(12)
        caption_para.font.color.rgb = self.colors['text']
        caption_para.alignment = PP_ALIGN.CENTER
    

    def save(self):
        timestamp = time.strftime("%Y%m%d-%H%M%S")
        version_filename = f"{self.filename.rsplit('.', 1)[0]}_{timestamp}.pptx"
        self.prs.save(version_filename)
        self.prs.save(self.filename)
        return version_filename
    
    def is_complete(self):
        return self.slide_count >= self.target_slides
    
    def get_progress(self):
        return f"Current progress: {self.slide_count}/{self.target_slides} slides"

# Agent configurations
content_strategist_config = {
    "name": "ContentStrategist",
    "system_message": """You are a presentation content strategist.
    Your responsibilities:
    - Create detailed presentation outlines
    - Define key messages and takeaways
    - Structure content flow and progression
    - Ensure content alignment with presentation goals
    - Recommend content distribution across slides
    
    For the AI trends presentation, focus on:
    1. Market size and growth
    2. Key technology developments
    3. Industry adoption rates
    4. Future predictions
    5. Challenges and opportunities"""
}

content_writer_config = {
    "name": "ContentWriter",
    "system_message": """You are a presentation content writer who MUST format ALL responses as valid JSON.

    CRITICAL REQUIREMENTS:
    1. Your ENTIRE response must be a single JSON object
    2. Use the EXACT format below - no deviations allowed:
    {
        "title": "Slide Title Here",
        "body": ["Bullet point 1", "Bullet point 2", "Bullet point 3"],
        "notes": "Speaker notes here"
    }

    3. Use double quotes for ALL strings
    4. Include exactly these three keys: title, body, notes
    5. The 'body' value must be an array of strings
    6. Do not include any explanatory text outside the JSON
    7. Do not use single quotes
    8. Do not use trailing commas
    9. Do not include any markdown formatting

    Example of CORRECT response:
    {
        "title": "Market Growth Analysis",
        "body": ["Global AI market reached $120B in 2024", "CAGR of 38% expected through 2028", "Healthcare sector leads adoption"],
        "notes": "Emphasize the rapid growth trajectory and highlight healthcare as key driver."
    }"""
}

visualization_expert_config = {
    "name": "VisualizationExpert",
    "system_message": """You are a data visualization specialist who MUST format ALL responses as valid JSON.

    CRITICAL REQUIREMENTS:
    1. Your ENTIRE response must be a single JSON object
    2. Use the EXACT format below - no deviations allowed:
    {
        "title": "Chart Title Here",
        "chart_type": "bar",
        "chart_data": {
            "labels": ["Label1", "Label2", "Label3"],
            "values": [10, 20, 30]
        }
    }

    Rules for chart_type:
    - Must be one of: "bar", "line", "pie"
    - Use lowercase only

    Rules for chart_data:
    - labels: array of strings
    - values: array of numbers
    - Both arrays must have the same length
    
    Example of CORRECT response:
    {
        "title": "AI Market Growth by Sector",
        "chart_type": "bar",
        "chart_data": {
            "labels": ["Healthcare", "Finance", "Retail"],
            "values": [45, 32, 28]
        }
    }

    Do not include any text outside the JSON structure."""
}

class CustomAssistantAgent(autogen.AssistantAgent):
    def __init__(self, name, system_message, presentation_builder):
        # Initialize parent class without the presentation_builder
        super().__init__(
            name=name,
            system_message=system_message,
            llm_config=False  # We'll use our own LLM implementation
        )
        
        # Store presentation_builder as instance variable
        self.presentation_builder = presentation_builder
        self.llm = OpenRouterLLM()
        
    def validate_json_structure(self, content):
        """Validate the JSON structure for content writer responses."""
        required_keys = {"title", "body", "notes"}
        
        if not all(key in content for key in required_keys):
            raise ValueError("Missing required keys in JSON response")
            
        if not isinstance(content["title"], str):
            raise ValueError("Title must be a string")
        if not isinstance(content["body"], list):
            raise ValueError("Body must be an array")
        if not all(isinstance(item, str) for item in content["body"]):
            raise ValueError("All body items must be strings")
        if not isinstance(content["notes"], str):
            raise ValueError("Notes must be a string")
            
        return True

    def validate_chart_structure(self, content):
        """Validate the JSON structure for visualization expert responses."""
        required_keys = {"title", "chart_type", "chart_data"}
        
        if not all(key in content for key in required_keys):
            raise ValueError("Missing required keys in JSON response")
        
        valid_chart_types = {"bar", "line", "pie"}
        if content["chart_type"] not in valid_chart_types:
            raise ValueError(f"Invalid chart type. Must be one of: {valid_chart_types}")
        
        required_data_keys = {"labels", "values"}
        if not all(key in content["chart_data"] for key in required_data_keys):
            raise ValueError("Chart data must contain 'labels' and 'values' arrays")
        
        if not isinstance(content["chart_data"]["labels"], list):
            raise ValueError("Labels must be an array")
        if not isinstance(content["chart_data"]["values"], list):
            raise ValueError("Values must be an array")
        if len(content["chart_data"]["labels"]) != len(content["chart_data"]["values"]):
            raise ValueError("Labels and values arrays must have the same length")
        
        if not all(isinstance(label, str) for label in content["chart_data"]["labels"]):
            raise ValueError("All labels must be strings")
        if not all(isinstance(value, (int, float)) for value in content["chart_data"]["values"]):
            raise ValueError("All values must be numbers")
        
        return True

    def format_response_as_json(self, text):
        """Attempt to extract and format JSON from the response text."""
        try:
            start = text.find('{')
            end = text.rfind('}') + 1
            
            if start >= 0 and end > start:
                json_str = text[start:end]
                return json.loads(json_str)
        except:
            # Default formatting as fallback
            lines = text.strip().split('\n')
            title = lines[0].strip()
            body = [line.strip().strip('*-').strip() for line in lines[1:-1] if line.strip()]
            notes = lines[-1].strip() if lines else "Speaker notes needed"
            
            return {
                "title": title,
                "body": body,
                "notes": notes
            }

    def format_chart_response_as_json(self, text):
        """Attempt to extract and format JSON from the visualization expert's response."""
        try:
            start = text.find('{')
            end = text.rfind('}') + 1
            
            if start >= 0 and end > start:
                json_str = text[start:end]
                return json.loads(json_str)
        except:
            # Default chart structure as fallback
            return {
                "title": "Data Visualization",
                "chart_type": "bar",
                "chart_data": {
                    "labels": ["Category A", "Category B", "Category C"],
                    "values": [30, 20, 10]
                }
            }

    def generate_reply(self, messages=None, sender=None, config=None):
        """Generate a reply based on the messages and handle slide creation."""
        if messages is None:
            messages = self._oai_messages
            
        if not messages:
            return None
            
        last_message = messages[-1]
        last_message_content = last_message.content if hasattr(last_message, 'content') else str(last_message)
            
        # Add formatting reminders based on agent type
        if self.name == "content_writer":
            last_message_content += "\n\nREMINDER: Your entire response must be a single valid JSON object with title, body, and notes keys."
        elif self.name == "visualization_expert":
            last_message_content += "\n\nREMINDER: Your entire response must be a single valid JSON object with title, chart_type, and chart_data keys."
            
        api_messages = [
            {"role": "system", "content": self.system_message},
            {"role": "user", "content": last_message_content}
        ]
        
        try:
            response = self.llm.create_completion(api_messages)
            reply = response['choices'][0]['message']['content']
            
            # Handle different agent types
            if self.name == "content_writer":
                try:
                    content_dict = self.format_response_as_json(reply)
                    self.validate_json_structure(content_dict)
                    self.presentation_builder.create_content_slide(content_dict)
                    print(f"\nAdded new slide: {content_dict['title']}")
                except (json.JSONDecodeError, ValueError) as e:
                    print(f"Warning: Content formatting error: {str(e)}")
                    default_content = {
                        "title": "Content Processing Error",
                        "body": [
                            "Error: Invalid content format received",
                            "This slide needs manual review",
                            f"Error details: {str(e)}"
                        ],
                        "notes": "This slide was created due to a content formatting error and needs to be reviewed."
                    }
                    self.presentation_builder.create_content_slide(default_content)
                    
            elif self.name == "visualization_expert":
                try:
                    chart_dict = self.format_chart_response_as_json(reply)
                    self.validate_chart_structure(chart_dict)
                    self.presentation_builder.create_chart_slide(
                        title=chart_dict['title'],
                        chart_type=chart_dict['chart_type'],
                        chart_data=chart_dict['chart_data']
                    )
                    print(f"\nAdded new chart slide: {chart_dict['title']}")
                except (json.JSONDecodeError, ValueError) as e:
                    print(f"Warning: Chart formatting error: {str(e)}")
                    default_chart = {
                        "title": "Visualization Error",
                        "chart_type": "bar",
                        "chart_data": {
                            "labels": ["Error", "Occurred"],
                            "values": [100, 0]
                        }
                    }
                    self.presentation_builder.create_chart_slide(**default_chart)
            
            return reply
            
        except Exception as e:
            print(f"Error generating reply: {str(e)}")
            return "I apologize, but I encountered an error processing your request."

def generate_presentation():
    # Initialize presentation builder
    builder = PresentationBuilder()
    
    # Create title slide
    builder.create_title_slide(
        "Artificial Intelligence Trends 2024",
        "A Comprehensive Market Analysis and Future Outlook"
    )
    
    # Initialize agents
    content_strategist = CustomAssistantAgent(
        name="content_strategist",
        system_message=content_strategist_config["system_message"],
        presentation_builder=builder
    )
    
    content_writer = CustomAssistantAgent(
        name="content_writer",
        system_message=content_writer_config["system_message"],
        presentation_builder=builder
    )
    
    visualization_expert = CustomAssistantAgent(
        name="visualization_expert",
        system_message=visualization_expert_config["system_message"],
        presentation_builder=builder
    )
    
    user_proxy = autogen.UserProxyAgent(
        name="user_proxy",
        human_input_mode="TERMINATE",
        max_consecutive_auto_reply=10
    )
    
    # Create group chat
    groupchat = autogen.GroupChat(
        agents=[user_proxy, content_strategist, content_writer, visualization_expert],
        messages=[],
        max_round=50
    )
    
    manager = autogen.GroupChatManager(groupchat=groupchat)
    
    # Generate presentation content
    initial_message = """Create a professional presentation about artificial intelligence trends in 2024.
    Requirements:
    1. Include executive summary
    2. Cover key trends, market analysis, and future predictions
    3. Include data visualizations and charts
    4. Provide speaker notes for each slide"""
    
    while not builder.is_complete():
        if builder.slide_count <= 1:  # Only title slide exists
            user_proxy.initiate_chat(manager, message=initial_message)
        else:
            next_slide_prompt = f"""Create the next slide content and design.
            {builder.get_progress()}
            Ensure content flows naturally from the previous slides."""
            user_proxy.initiate_chat(manager, message=next_slide_prompt)
    
    # Save final presentation
    final_file = builder.save()
    print(f"\nPresentation completed: {final_file}")
    print(builder.get_progress())
    
    return final_file

if __name__ == "__main__":
    
    # Generate the presentation
    output_file = generate_presentation()
    print(f"Presentation generated successfully: {output_file}")

